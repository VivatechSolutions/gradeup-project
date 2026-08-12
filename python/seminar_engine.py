"""
AI Seminar Engine for GradeUp

Two session modes:
  - DEMO / PRACTICE: AI coaches the student, guides them through the seminar,
    provides RAG-based guidance on demand via guide_student().
    Storage: seminar_practice_data/
  - MAIN: AI examines the student — no guidance, only assessment.
    Storage: seminar_data/

Common features:
- Student selects unit, AI presents important topics
- Student starts presenting, AI listens and interacts
- If student pauses 8-10 seconds → AI gives contextual hints
- AI can repeat, clarify, rephrase when student asks
- Scoring: Conceptual Understanding, Depth, Flow, Engagement, Hints Used
"""

import os
import json
import hashlib
import time
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime, timezone

import requests

SEMINAR_DATA_DIR = Path("seminar_data")
SEMINAR_PRACTICE_DATA_DIR = Path("seminar_practice_data")
SEMINAR_PRACTICE_SESSION_DIR = Path("practice_session")
SEMINAR_CHAT_DIR = Path("seminar_chat")
SEMINAR_MODEL = "gpt-4o-mini"
SEMINAR_FALLBACK_MODEL = "gpt-4o"
SEMINAR_TIMEOUT = 90
HINT_SILENCE_THRESHOLD = 8  # seconds of silence to trigger hint


# ── File Content Extraction Utilities ─────────────────────────────────────────

def extract_uploaded_file_content(filename: str, file_bytes: bytes) -> str:
    """
    Extract text content from an uploaded PDF or PPTX file.

    Uses PyMuPDF (fitz) for PDF and python-pptx for PPTX.
    Returns the extracted text as a single string.
    """
    fname_lower = filename.lower()

    if fname_lower.endswith(".pdf"):
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages_text = []
            for page in doc:
                pages_text.append(page.get_text())
            doc.close()
            text = "\n\n".join(pages_text)
            print(f"  📄 [SeminarEngine] Extracted {len(text):,} chars from PDF ({len(pages_text)} pages)")
            return text
        except ImportError:
            raise ImportError("PyMuPDF (fitz) is required for PDF extraction. Install: pip install PyMuPDF")

    elif fname_lower.endswith(".pptx"):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_parts.append(para_text)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                slide_parts.append(row_text)
                slides_text.append("\n".join(slide_parts))
            text = "\n\n".join(slides_text)
            print(f"  📊 [SeminarEngine] Extracted {len(text):,} chars from PPTX ({len(slides_text)} slides)")
            return text
        except ImportError:
            raise ImportError("python-pptx is required for PPTX extraction. Install: pip install python-pptx")

    elif fname_lower.endswith(".ppt"):
        # Legacy .ppt format — limited support
        raise ValueError(
            "Legacy .ppt format is not supported. Please convert to .pptx format."
        )

    else:
        raise ValueError(f"Unsupported file type: {filename}. Only PDF and PPTX are supported.")


class SeminarEngine:
    """
    Manages AI Seminar sessions with interactive features,
    hint system, topic navigation, and student assessment.
    """

    def __init__(
        self,
        data_dir: Path = SEMINAR_DATA_DIR,
        practice_data_dir: Path = SEMINAR_PRACTICE_DATA_DIR,
        practice_session_dir: Path = SEMINAR_PRACTICE_SESSION_DIR,
        chat_dir: Path = SEMINAR_CHAT_DIR,
    ):
        self.data_dir = data_dir
        self.practice_data_dir = practice_data_dir
        self.practice_session_dir = practice_session_dir
        self.chat_dir = chat_dir
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.practice_data_dir.mkdir(parents=True, exist_ok=True)
        self.practice_session_dir.mkdir(parents=True, exist_ok=True)
        self.chat_dir.mkdir(parents=True, exist_ok=True)

    def _session_path(self, session_id: str, session_mode: str = "main") -> Path:
        """Return path for a session file. Demo → practice_data, main → data, practice → practice_session."""
        if session_mode == "practice":
            base = self.practice_session_dir
        elif session_mode == "demo":
            base = self.practice_data_dir
        else:
            base = self.data_dir
        return base / f"seminar__{session_id}.json"

    def _load_session(self, session_id: str) -> Optional[Dict[str, Any]]:
        """Load session — checks main, practice_data, and practice_session directories."""
        for base_dir in (self.data_dir, self.practice_data_dir, self.practice_session_dir):
            path = base_dir / f"seminar__{session_id}.json"
            if path.exists():
                try:
                    return json.loads(path.read_text(encoding="utf-8"))
                except (json.JSONDecodeError, OSError):
                    pass
        return None

    def _save_session(self, session: Dict[str, Any]) -> None:
        mode = session.get("session_mode", "main")
        path = self._session_path(session["session_id"], mode)
        session["updated_at"] = datetime.now(timezone.utc).isoformat()
        path.write_text(
            json.dumps(session, indent=2, ensure_ascii=False), encoding="utf-8"
        )

    def _generate_session_id(self, candidate_id: str) -> str:
        seed = f"seminar_{candidate_id}_{datetime.now().isoformat()}_{time.time()}"
        return hashlib.md5(seed.encode()).hexdigest()[:16]

    def _call_llm(self, messages: List[Dict], temperature: float = 0.7) -> str:
        api_key = os.environ.get("OPENAI_API_KEY_TEXT") or os.environ.get("OPENAI_API_KEY")
        if not api_key:
            return "AI is not configured."

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": SEMINAR_MODEL,
            "messages": messages,
            "max_completion_tokens": 2048,
            "temperature": temperature,
        }

        try:
            resp = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers=headers, json=payload, timeout=SEMINAR_TIMEOUT,
            )
            if not resp.ok:
                payload["model"] = SEMINAR_FALLBACK_MODEL
                resp = requests.post(
                    "https://api.openai.com/v1/chat/completions",
                    headers=headers, json=payload, timeout=SEMINAR_TIMEOUT,
                )
            if resp.ok:
                return resp.json()["choices"][0]["message"]["content"].strip()
        except Exception as e:
            print(f"  ❌ [SeminarEngine] Error: {e}")
        return "Something went wrong. Please try again."

    def _call_llm_json(self, messages: List[Dict], temperature: float = 0.3) -> Optional[Dict]:
        raw = self._call_llm(messages, temperature)
        try:
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            start = raw.find("{")
            end = raw.rfind("}")
            if start != -1 and end != -1:
                return json.loads(raw[start:end + 1])
        except Exception:
            pass
        return None

    # ── System Prompts ────────────────────────────────────────────────────────

    def _seminar_system_prompt(
        self,
        subject: str,
        unit_number: int,
        current_topic: str,
        candidate_name: str,
        rag_context: str,
    ) -> str:
        return f"""You are GradeUp AI Seminar Examiner — an interactive, patient, and encouraging academic examiner for school students.

## YOUR ROLE
You are examining {candidate_name}'s knowledge on "{current_topic}" from {subject}, Unit {unit_number}.
You are like a friendly examiner in a viva/oral exam — listen carefully, ask follow-up questions, and assess understanding.

## INTERACTION RULES
1. **Be HIGHLY interactive**: Ask follow-up questions, validate good explanations, correct misconceptions.
2. **Repeat when asked**: If the student says "I didn't understand", "repeat", "say again", or similar → Rephrase your last point in simpler terms.
3. **Clarify**: If the student seems confused, provide a simpler explanation using everyday analogies.
4. **Encourage**: Praise good explanations: "Excellent point!", "That's a great way to explain it!"
5. **Correct gently**: If the student is wrong, say "That's close, but..." or "Actually, according to the textbook..."
6. **Stay focused**: Keep the discussion on the current topic. If the student drifts, gently redirect.
7. **Keep responses SHORT**: 2-3 sentences max. This is a conversation, not a lecture.

## CONTENT SAFETY
- NEVER produce inappropriate, 18+, or harmful content.
- If student sends inappropriate content, redirect: "Let's focus on our topic! Tell me about..."

## TEXTBOOK CONTEXT (use this to verify student's explanations and provide corrections)
{rag_context if rag_context else "(No context available)"}

## RESPONSE STYLE
- Conversational and warm
- Address {candidate_name} by name sometimes
- Use questions to probe understanding: "Can you explain why?", "What do you think would happen if...?"
- Validate before correcting: "Yes, and also..." or "Good start, but consider..."
"""

    def _demo_system_prompt(
        self,
        subject: str,
        unit_number: int,
        current_topic: str,
        candidate_name: str,
        rag_context: str,
    ) -> str:
        """System prompt for DEMO / PRACTICE sessions — AI acts as a coach."""
        return f"""You are GradeUp AI Seminar Coach — a friendly, supportive coach helping school students practice and prepare for their seminar presentations.

## YOUR ROLE
You are coaching {candidate_name} to practice presenting on "{current_topic}" from {subject}, Unit {unit_number}.
You are NOT an examiner — you are a helpful guide. Your job is to help the student learn how to present well.

## COACHING RULES
1. **Guide actively**: Tell the student what to cover, suggest how to structure their explanation, and give tips on presentation flow.
2. **Use the textbook content**: Reference the textbook context below to help the student cover all important points.
3. **Correct with full explanations**: If the student is wrong, don't just say "that's close" — explain the correct concept fully so they learn.
4. **Suggest what to say next**: If the student finishes a point, suggest what they should cover next based on the topic.
5. **Teach seminar skills**: Offer tips like "Start with a definition", "Use an example here", "Summarize your key points", "Connect this to real life".
6. **Encourage and build confidence**: This is practice — be positive and supportive. "Great start!", "You're improving!", "That's exactly the right approach!"
7. **Response Length**: During the actual presentation coaching, keep responses to 3-5 sentences max. However, for the initial introduction, you MUST provide a detailed and comprehensive briefing without any length constraints. Give direction without lecturing.

## CONTENT SAFETY
- NEVER produce inappropriate, 18+, or harmful content.
- If student sends inappropriate content, redirect: "Let's focus on practicing! Try explaining..."

## TEXTBOOK CONTEXT (use this to guide the student's content)
{rag_context if rag_context else "(No context available)"}

## RESPONSE STYLE
- Warm, encouraging, and directive (like a good coach)
- Address {candidate_name} by name sometimes
- Use phrases like: "Now try explaining...", "A good point to mention here would be...", "You could strengthen this by adding..."
- After the student explains something, validate what was good AND suggest improvements
"""

    # ── Practice Step Definitions ──────────────────────────────────────────

    PRACTICE_STEPS = {
        1: {"name": "Topic Understanding", "description": "Understand the core concepts of the topic"},
        2: {"name": "Slide Planning", "description": "Plan the slide structure and content"},
        3: {"name": "Tech Guidance", "description": "Learn how to create the PDF/PPT using your chosen tool"},
        4: {"name": "Presentation Practice", "description": "Practice presenting each slide verbally"},
        5: {"name": "Follow-up Q&A", "description": "Answer follow-up questions to test readiness"},
    }

    def _practice_system_prompt(
        self,
        subject: str,
        unit_number: int,
        current_topic: str,
        candidate_name: str,
        rag_context: str,
        practice_step: int = 1,
        chosen_tool: str = "",
    ) -> str:
        """System prompt for PRACTICE sessions — step-aware, interactive, English-only, with tool-specific guidance."""

        step_info = self.PRACTICE_STEPS.get(practice_step, self.PRACTICE_STEPS[1])
        step_name = step_info["name"]
        step_desc = step_info["description"]

        # Tool-specific guidance block
        tool_guidance = ""
        if practice_step == 3:
            if chosen_tool:
                tool_lower = chosen_tool.lower()
                if "google" in tool_lower or "slides" in tool_lower:
                    tool_guidance = """
## TOOL-SPECIFIC GUIDANCE: GOOGLE SLIDES
1. Go to slides.google.com and sign in with your Google account.
2. Click "+ Blank presentation" or choose a template.
3. First slide: Click the title text box → type your topic name. Add subtitle with your name, subject, unit, date.
4. Add new slides: Click "+ New Slide" in the toolbar (or Ctrl+M).
5. For each content slide: Use the "Title and Body" layout. Type the heading, then add bullet points.
6. To add images: Insert → Image → Search the web or Upload from computer.
7. Design tips: Use Theme (Slide → Change theme) for consistent colors. Keep text large (18pt+ body, 24pt+ headings).
8. To export as PDF: File → Download → PDF Document (.pdf).
9. To share: File → Share → copy the link.
"""
                elif "canva" in tool_lower:
                    tool_guidance = """
## TOOL-SPECIFIC GUIDANCE: CANVA
1. Go to canva.com and sign up or log in (free account works).
2. Click "Create a design" → search "Presentation" → pick a template.
3. First slide: Click on text to edit. Add your topic name, your name, subject, unit, and date.
4. Add new slides: Click "+ Add page" at the bottom.
5. For content slides: Click text boxes to edit. Use the left panel to drag in elements, icons, or images.
6. To add images: Click "Uploads" in left panel → upload your images, or use "Elements" → search for graphics.
7. Design tips: Stick to 2-3 colors. Use Canva's built-in font combinations. Keep each slide clean.
8. To export as PDF: Click "Share" (top right) → "Download" → choose "PDF Standard" → Download.
9. To present directly: Click "Present" button in top right.
"""
                elif "powerpoint" in tool_lower or "ppt" in tool_lower or "microsoft" in tool_lower:
                    tool_guidance = """
## TOOL-SPECIFIC GUIDANCE: MICROSOFT POWERPOINT
1. Open PowerPoint on your computer (or use PowerPoint Online at office.com).
2. Choose a template or start with a blank presentation.
3. First slide: Click "Click to add title" → type your topic name. Add subtitle with your name, subject, unit, date.
4. Add new slides: Home tab → New Slide → choose a layout (use "Title and Content" for most slides).
5. For each content slide: Type the heading in the title area, then add bullet points in the content area.
6. To add images: Insert tab → Pictures → choose from your computer or online.
7. Design tips: Use the Design tab to pick a consistent theme. Keep text at 18pt+ for body, 28pt+ for titles.
8. To export as PDF: File → Save As → choose "PDF" from the file type dropdown. Or File → Export → Create PDF.
9. Keyboard shortcuts: Ctrl+M (new slide), Ctrl+S (save), F5 (start slideshow).
"""
                elif "libre" in tool_lower or "impress" in tool_lower:
                    tool_guidance = """
## TOOL-SPECIFIC GUIDANCE: LIBREOFFICE IMPRESS
1. Open LibreOffice Impress (free download from libreoffice.org).
2. Choose a template from the startup wizard or start blank.
3. First slide: Click on the title placeholder → type your topic name. Add your name, subject, unit, and date.
4. Add new slides: Right-click in the slide panel (left side) → Insert New Slide.
5. For content slides: Use "Blank" or "Title, Content" layouts. Add text boxes via Insert → Text Box.
6. To add images: Insert → Image → choose your file.
7. Design tips: Use Slide → Slide Properties to set consistent backgrounds. Keep fonts readable (18pt+ body).
8. To export as PDF: File → Export as PDF → click "Export" → save.
"""
                else:
                    tool_guidance = f"""
## TOOL-SPECIFIC GUIDANCE: {chosen_tool.upper()}
Provide step-by-step guidance for using {chosen_tool} to create a presentation. Cover: creating slides, adding text and images, choosing templates, and exporting as PDF. If you are not sure about the exact steps for this tool, provide general guidance that applies to most presentation tools.
"""
            else:
                tool_guidance = """
## TOOL SELECTION NEEDED
You MUST ask the student which tool they want to use to create their presentation. Ask:
"Which tool would you like to use to create your slides? Here are some popular options:
1. 📊 Google Slides (free, works in browser)
2. 🎨 Canva (free, great templates)
3. 💼 Microsoft PowerPoint (if you have it installed)
4. 🆓 LibreOffice Impress (free, open source)
5. Or tell me any other tool you prefer!"

Wait for their answer before giving specific instructions.
"""

        return f"""You are GradeUp AI Seminar Preparation Coach — an interactive, step-by-step mentor helping school students learn HOW to prepare and deliver a great seminar presentation.

## YOUR ROLE
You are coaching {candidate_name} on preparing a seminar on "{current_topic}" from {subject}, Unit {unit_number}.
You are NOT an examiner. You are a coach. Be interactive — every response MUST end with a question or action item.

## LANGUAGE
Always respond in **English only**, regardless of what language the student types in.

## CURRENT STEP: Step {practice_step} — {step_name}
**Goal**: {step_desc}

The practice session has 5 steps:
1. Topic Understanding — Check what the student knows, fill gaps
2. Slide Planning — Plan slide structure and content together
3. Tech Guidance — Teach how to create PDF/PPT using their chosen tool
4. Presentation Practice — Student practices presenting verbally
5. Follow-up Q&A — Ask 3-5 questions to test readiness

## STEP-SPECIFIC INSTRUCTIONS

### If Step 1 (Topic Understanding):
- Ask the student to explain what they know about "{current_topic}" in their own words.
- Listen to their response, correct any mistakes using the textbook context.
- Fill in important concepts they missed.
- When they demonstrate basic understanding of the key concepts, tell them: "Great! You have a good grasp of the topic. Let's move to planning your slides!" and transition to Step 2.
- If the student explicitly asks to skip ahead or move to slides/tech/practice, respect their request.

### If Step 2 (Slide Planning):
- Guide them to plan a 5-8 slide presentation:
  - Slide 1: Title (topic, name, subject, unit, date)
  - Slide 2: Introduction/Overview
  - Slides 3-5: Key concepts (one per slide, bullet points)
  - Slide 6: Examples/Real-world applications
  - Slide 7: Summary/Conclusion
- Ask them to describe what they'd put on each slide. Give feedback.
- Key rule: "Slides are your OUTLINE — you explain MORE than what's on the slide."
- When they have a solid slide plan, transition to Step 3.
- If the student asks to skip to tech guidance or practice, respect their request.

### If Step 3 (Tech Guidance):
{tool_guidance}
- Walk them through the tool step-by-step.
- Offer design tips: font sizes (18pt+ body, 24pt+ headings), 3-5 bullet points per slide, consistent colors, minimal text.
- Common mistakes: wall of text, too many animations, reading from slides, cramming content.
- When they feel confident about creating slides, transition to Step 4.
- If the student asks to skip to practice or Q&A, respect their request.

### If Step 4 (Presentation Practice):
- Ask the student to practice presenting one slide at a time.
- After each attempt, give structured feedback:
  - ✅ What you did well
  - 🔧 What to improve
  - 💡 Try this next
- Coach delivery: "Start with a definition", "Use an example", "Add a transition phrase like 'This leads us to...'"
- When they've practiced the key slides, transition to Step 5.
- If the student asks to skip to Q&A, respect their request.

### If Step 5 (Follow-up Q&A):
- Ask 3-5 questions about "{current_topic}" to test the student's readiness.
- Ask ONE question at a time. Wait for their answer.
- After each answer: validate what's correct, gently correct mistakes, explain the right answer.
- After all questions, give a final summary: "You're ready for your seminar! Here's a quick recap of what to remember..."
- Tell them they can now go to the Demo or Main session to present.

## INTERACTION RULES
1. **ALWAYS end with a question or action item** — never leave the student without a next step.
2. Keep responses to 3-6 sentences (except for initial intro or detailed tech guidance).
3. Be warm, encouraging, and specific.
4. Correct mistakes fully using the textbook context.
5. Use phrases like: "Now try...", "What would you say about...", "Can you explain...?"
6. Track progress — when the student is ready, explicitly say you're moving to the next step.

## STEP TRANSITION
When you determine the student is ready to move to the next step, include this EXACT marker in your response:
[STEP_COMPLETE:{practice_step}]
This signals the system to advance the step counter. Only include this when the student has genuinely completed the current step's goals OR when they explicitly ask to move ahead.

## CONTENT SAFETY
- NEVER produce inappropriate, 18+, or harmful content.
- If student sends inappropriate content, redirect: "Let's focus on preparing your seminar!"

## TEXTBOOK CONTEXT
{rag_context if rag_context else "(No context available)"}

## RESPONSE STYLE
- Warm, encouraging, and interactive
- Address {candidate_name} by name sometimes
- Always end with a follow-up question or action item
- After practice attempts, use structured feedback format

## SUGGESTED FOLLOW-UP QUESTIONS
At the END of EVERY response, you MUST include exactly 2-3 short suggested questions/responses that the student can use as quick replies. Format them EXACTLY like this (on a new line at the very end):
[SUGGESTIONS: "suggested reply 1" | "suggested reply 2" | "suggested reply 3"]

Examples based on step:
- Step 1: [SUGGESTIONS: "Distance is the total path length" | "I'm not sure, can you explain?" | "Let me try explaining"]
- Step 2: [SUGGESTIONS: "I'll start with the title slide" | "What should go on slide 3?" | "How many slides do I need?"]
- Step 3: [SUGGESTIONS: "I'll use Google Slides" | "I'll use Canva" | "I'll use PowerPoint"]
- Step 4: [SUGGESTIONS: "Let me try presenting slide 1" | "Can you show me an example?" | "What should I say first?"]
- Step 5: [SUGGESTIONS: "I think the answer is..." | "Can you give me a hint?" | "I'm not sure about this one"]

The suggestions should be contextually relevant to what you just asked. NEVER skip this.
"""

    def _compare_rag_vs_uploaded(
        self,
        rag_context: str,
        uploaded_content: str,
        topic: str,
        subject: str,
    ) -> Dict[str, Any]:
        """
        Compare RAG-retrieved textbook content with student-uploaded content.

        Identifies:
          - Topics present in uploaded file but missing from RAG
          - Topics present in RAG but missing from uploaded file
          - Common topics covered by both
          - An overall analysis summary
        """
        prompt = f"""You are an educational content analyst. Compare the following two sources of content about the topic "{topic}" from {subject}.

## SOURCE 1: TEXTBOOK (RAG-retrieved)
{rag_context[:5000]}

## SOURCE 2: STUDENT UPLOADED FILE
{uploaded_content[:5000]}

## YOUR TASK
Analyze both sources and provide a structured comparison. Identify:
1. Topics/concepts present in the UPLOADED FILE but MISSING from the TEXTBOOK (RAG)
2. Topics/concepts present in the TEXTBOOK (RAG) but MISSING from the UPLOADED FILE
3. Common topics covered by BOTH sources
4. An overall analysis summary (2-3 sentences)

Return ONLY valid JSON in this format:
{{
    "missing_in_upload": ["topic/concept not in uploaded file but in textbook", ...],
    "analysis_summary": "Brief 2-3 sentence overall comparison.",
    "uploaded_content_quality": "good|adequate|poor",
    "recommendation": "Brief recommendation for the student about content gaps."
}}"""

        raw = self._call_llm([
            {"role": "system", "content": "You are a precise educational content analyst. Always respond with valid JSON only."},
            {"role": "user", "content": prompt},
        ], temperature=0.3)

        # Parse the response
        try:
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            start = raw.find("{")
            end = raw.rfind("}")
            if start != -1 and end != -1:
                result = json.loads(raw[start:end + 1])
                return result
        except Exception as e:
            print(f"  ⚠️ [SeminarEngine] Failed to parse comparison analysis: {e}")

        # Fallback
        return {
            "missing_in_upload": [],
            "analysis_summary": "Comparison analysis could not be completed.",
            "uploaded_content_quality": "unknown",
            "recommendation": "Please review both sources manually.",
        }

    def _get_cross_session_history(
        self,
        candidate_id: str,
        subject: str,
        unit_number: int,
        topic: str,
    ) -> Optional[Dict[str, Any]]:
        """
        Check if the student has completed this topic in a debate session.
        Returns the latest debate session info (score, session_id) or None.
        """
        debate_dir = Path("debate_data")
        if not debate_dir.exists():
            return None

        best_match = None
        for path in debate_dir.glob("debate__*.json"):
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
                if (
                    data.get("candidate_id") == candidate_id
                    and (data.get("subject") or "").lower() == subject.lower()
                    and data.get("unit_number") == unit_number
                    and (data.get("topic") or "").lower() == topic.lower()
                    and data.get("status") == "ended"
                    and data.get("scores")
                ):
                    score = data["scores"].get("total_score", 0)
                    created = data.get("created_at", "")
                    if best_match is None or created > best_match["created_at"]:
                        best_match = {
                            "session_id": data.get("session_id"),
                            "topic": data.get("topic"),
                            "score": score,
                            "created_at": created,
                            "session_type": "debate",
                        }
            except (json.JSONDecodeError, OSError):
                continue
        return best_match

    # ── Helper: Extract suggested questions from AI response ─────────────────

    @staticmethod
    def _extract_suggested_questions(ai_response: str) -> tuple:
        """
        Extract [SUGGESTIONS: ...] marker from AI response.
        Returns (cleaned_response, list_of_suggestions).
        """
        import re as _re
        match = _re.search(r'\[SUGGESTIONS:\s*(.+?)\]', ai_response, _re.DOTALL)
        if match:
            raw = match.group(1)
            suggestions = [s.strip().strip('"').strip("'") for s in raw.split("|")]
            suggestions = [s for s in suggestions if s]  # filter empty
            cleaned = _re.sub(r'\[SUGGESTIONS:\s*.+?\]', '', ai_response, flags=_re.DOTALL).strip()
            return cleaned, suggestions
        return ai_response, []


    # ── Public API ────────────────────────────────────────────────────────────

    def start_seminar(
        self,
        candidate_id: str,
        candidate_name: str,
        subject: str,
        unit_number: int,
        board: str,
        class_number: str,
        unit_name: str = "",
        topic: Optional[str] = None,
        session_mode: str = "main",
        uploaded_content: Optional[str] = None,
        term: Optional[Any] = None,
    ) -> Dict[str, Any]:
        """
        Start a new seminar session.

        session_mode:
          - "demo"     → Practice/rehearsal with AI coaching + hints
          - "main"     → Real exam with AI examining (default)
          - "practice" → Interactive seminar preparation training

        Topic is mandatory — student must select it. Validated against RAG.
        """
        import re
        from ai_tutor import retrieve_context, _format_context

        if session_mode not in ("demo", "main", "practice"):
            return {"error": "session_mode must be 'demo', 'main', or 'practice'", "success": False}

        if not topic or not topic.strip():
            return {"error": "topic is required — students must select their seminar topic", "success": False}

        # ── Enforce mandatory file upload for main/demo ──────────────────
        if session_mode in ("main", "demo") and not uploaded_content:
            return {
                "error": "PDF or PPT file upload is mandatory for demo and main seminar sessions. "
                         "Please prepare and upload your presentation material before starting.",
                "success": False,
            }

        current_topic = re.sub(r'^[\d.]+\s+', '', topic.strip())

        session_id = self._generate_session_id(candidate_id)

        # Get RAG context and validate topic
        chunks = retrieve_context(
            query=current_topic, subject=subject, unit_number=unit_number,
            board=board, class_number=class_number, limit=5, term=term
        )
        rag_context = _format_context(chunks)

        if not rag_context and not uploaded_content:
            return {"error": f"topic '{current_topic}' not found in course material", "success": False}

        # ── Determine session content source ───────────────────────────────
        # If student uploaded a file (main/demo only), use that as primary context
        content_source = "rag"  # default
        comparison_analysis = None
        session_context = rag_context  # default to RAG

        if uploaded_content and session_mode in ("main", "demo"):
            content_source = "uploaded_file"
            session_context = uploaded_content
            print(f"  📎 [SeminarEngine] Using uploaded file content as session source ({len(uploaded_content):,} chars)")
        elif not rag_context and uploaded_content:
            # RAG returned nothing but file was provided
            content_source = "uploaded_file"
            session_context = uploaded_content
            print(f"  📎 [SeminarEngine] RAG empty — using uploaded file as fallback")

        # ── Generate mode-appropriate greeting ─────────────────────────────
        if session_mode == "practice":
            system_prompt = self._practice_system_prompt(
                subject, unit_number, current_topic, candidate_name, rag_context,
                practice_step=1,
            )
            greeting_prompt = f"""Welcome {candidate_name} to their SEMINAR PRACTICE SESSION on "{current_topic}" from {subject}, Unit {unit_number}.

This is Step 1 — Topic Understanding. Keep your introduction SHORT and interactive (under 6 sentences total):

1. Greet {candidate_name} warmly (1 sentence). Tell them this is a safe, no-pressure practice space.
2. Briefly mention: "We'll go through 5 steps together — understanding the topic, planning your slides, creating your PDF/PPT, practicing your presentation, and some Q&A." (1 sentence)
3. Mention the 2-3 most important concepts from "{current_topic}" that they should know about (based on textbook context). Keep it brief — just name them. (1-2 sentences)
4. Ask your first interactive question to start Step 1: "Let's start! Can you tell me in your own words what {current_topic} is about?" (1 sentence)

Do NOT give a long lecture. Do NOT explain how to make slides yet. Just greet, set context, and ask the first question.

At the END of your response, include 2-3 suggested quick replies in this EXACT format:
[SUGGESTIONS: "suggested reply 1" | "suggested reply 2" | "suggested reply 3"]
Make them relevant to the first question you ask."""
        elif session_mode == "demo":
            system_prompt = self._demo_system_prompt(
                subject, unit_number, current_topic, candidate_name, rag_context
            )
            greeting_prompt = f"""Welcome {candidate_name} to their PRACTICE seminar session on "{current_topic}" from {subject}, Unit {unit_number}.

Keep your introduction SHORT and actionable (under 10 sentences total):
1. Greet {candidate_name} warmly (1 sentence). Mention this is a safe practice space.
2. Briefly explain the seminar format: "You will present, I will listen and guide you. After your presentation, I'll ask follow-up questions." (2 sentences max)
3. Give a BRIEF topic briefing on "{current_topic}" — mention the 3-4 key concepts they should aim to cover based on the textbook context. (3-4 sentences)
{'4. Note: The session is based on the students uploaded material. Mention that their uploaded content will be the primary source.' if content_source == 'uploaded_file' else ''}
{'5. ' if content_source == 'uploaded_file' else '4. '}Invite them to start when ready. (1 sentence)

Be concise and practical — do NOT give a long lecture. Save the detailed coaching for when they start presenting."""
        else:
            system_prompt = self._seminar_system_prompt(
                subject, unit_number, current_topic, candidate_name, session_context
            )
            uploaded_note = ""
            if content_source == "uploaded_file":
                uploaded_note = "\nNote: This session is based on the student's uploaded material (PDF/presentation), not the textbook. Reference their uploaded content."
            greeting_prompt = f"""Greet {candidate_name} for their seminar session on {subject}, Unit {unit_number}.
Present the topic: "{current_topic}".
Provide a brief overview of the topic — mention the 3-4 key concepts or areas the student should aim to cover in their presentation (based on the context provided).
Invite the student to start explaining this topic.
Tell them you'll be listening and asking questions to test their understanding.
Keep the greeting warm and encouraging. Limit to 6-8 sentences.{uploaded_note}"""

        ai_greeting = self._call_llm([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": greeting_prompt},
        ])

        session = {
            "session_id": session_id,
            "candidate_id": candidate_id,
            "candidate_name": candidate_name,
            "subject": subject,
            "unit_number": unit_number,
            "unit_name": unit_name,
            "board": board,
            "class_number": class_number,
            "current_topic": current_topic,
            "rag_context": rag_context,
            "session_context": session_context,
            "content_source": content_source,
            "session_mode": session_mode,
            "status": "active",
            "messages": [
                {
                    "role": "ai",
                    "content": ai_greeting,
                    "timestamp": datetime.now(timezone.utc).isoformat(),
                    "type": "greeting",
                }
            ],
            "hints_given": 0,
            "total_interactions": 0,
            "scores": None,
            "comparison_analysis": comparison_analysis,
            "created_at": datetime.now(timezone.utc).isoformat(),
        }

        # ── Practice session: add step tracking ───────────────────────────
        if session_mode == "practice":
            session["practice_step"] = 1
            session["practice_step_name"] = "Topic Understanding"
            session["chosen_tool"] = ""
            session["qa_questions_asked"] = 0

        self._save_session(session)

        result = {
            "session_id": session_id,
            "session_mode": session_mode,
            "current_topic": current_topic,
            "ai_greeting": ai_greeting,
            "content_source": content_source,
        }

        # Include practice step info
        if session_mode == "practice":
            # Extract suggested questions from greeting
            cleaned_greeting, suggestions = self._extract_suggested_questions(ai_greeting)
            result["ai_greeting"] = cleaned_greeting
            result["practice_step"] = 1
            result["practice_step_name"] = "Topic Understanding"
            result["practice_progress"] = 0
            result["suggested_questions"] = suggestions

        return result

    def seminar_respond(
        self,
        session_id: str,
        student_message: str,
        silence_seconds: float = 0,
    ) -> Dict[str, Any]:
        """
        Handle the student's seminar presentation input.

        Behavior varies by session_mode:
          - "main":     Record only — no AI response generated.
          - "demo":     Record + return a short AI hint (1-2 sentence nudge).
          - "practice": Record + return a full interactive AI coaching response.
        """
        session = self._load_session(session_id)
        if not session:
            return {"error": "Session not found", "success": False}
        if session.get("status") == "ended":
            return {"error": "Session already ended", "success": False}

        session_mode = session.get("session_mode", "main")

        # ── Record student message ────────────────────────────────────────
        session["messages"].append({
            "role": "student",
            "content": student_message,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })
        session["total_interactions"] = session.get("total_interactions", 0) + 1

        student_turns = sum(1 for m in session["messages"] if m.get("role") == "student")

        # ── Mode-specific AI response ─────────────────────────────────────
        if session_mode == "demo":
            # Demo: return a short hint (same style as silence-triggered hint)
            hint_result = self._give_hint(session)
            return {
                "success": True,
                "current_topic": session["current_topic"],
                "turn_number": student_turns,
                "hints_used": session.get("hints_given", 0),
                "session_id": session_id,
                "ai_response": hint_result.get("ai_response", ""),
                "is_hint": True,
            }

        elif session_mode == "practice":
            # Practice: full interactive AI coaching with step tracking
            rag_context = session.get("rag_context", "")
            practice_step = session.get("practice_step", 1)
            chosen_tool = session.get("chosen_tool", "")

            # ── Detect tool choice from student message ───────────────────
            if practice_step == 3 and not chosen_tool:
                msg_lower = student_message.lower()
                tool_keywords = {
                    "google slides": ["google", "slides", "google slides"],
                    "canva": ["canva"],
                    "microsoft powerpoint": ["powerpoint", "ppt", "microsoft", "ms powerpoint"],
                    "libreoffice impress": ["libreoffice", "libre", "impress"],
                }
                for tool_name, keywords in tool_keywords.items():
                    if any(kw in msg_lower for kw in keywords):
                        chosen_tool = tool_name
                        session["chosen_tool"] = chosen_tool
                        print(f"  🔧 [SeminarEngine] Practice: Student chose tool: {chosen_tool}")
                        break
                # If no match but step 3, treat the whole message as tool name
                if not chosen_tool and len(student_message.strip()) < 50:
                    chosen_tool = student_message.strip()
                    session["chosen_tool"] = chosen_tool
                    print(f"  🔧 [SeminarEngine] Practice: Student chose tool (raw): {chosen_tool}")

            system_prompt = self._practice_system_prompt(
                session["subject"], session["unit_number"],
                session["current_topic"], session["candidate_name"],
                rag_context,
                practice_step=practice_step,
                chosen_tool=chosen_tool,
            )

            # Build conversation history for LLM
            llm_messages = [{"role": "system", "content": system_prompt}]
            for msg in session["messages"][-20:]:
                role = "assistant" if msg["role"] == "ai" else "user"
                llm_messages.append({"role": role, "content": msg["content"]})

            ai_response = self._call_llm(llm_messages)

            # ── Detect step completion marker ─────────────────────────────
            import re as _re
            step_complete_match = _re.search(r'\[STEP_COMPLETE:(\d+)\]', ai_response)
            if step_complete_match:
                completed_step = int(step_complete_match.group(1))
                if completed_step == practice_step and practice_step < 5:
                    practice_step += 1
                    session["practice_step"] = practice_step
                    step_info = self.PRACTICE_STEPS.get(practice_step, {})
                    session["practice_step_name"] = step_info.get("name", "")
                    print(f"  📈 [SeminarEngine] Practice: Advanced to Step {practice_step} — {step_info.get('name', '')}")
                # Remove the marker from the response shown to student
                ai_response = _re.sub(r'\[STEP_COMPLETE:\d+\]', '', ai_response).strip()

            # ── Track Q&A questions in Step 5 ─────────────────────────────
            if practice_step == 5:
                session["qa_questions_asked"] = session.get("qa_questions_asked", 0) + 1

            # ── Extract suggested questions from AI response ───────────
            ai_response, suggested_questions = self._extract_suggested_questions(ai_response)

            session["messages"].append({
                "role": "ai",
                "content": ai_response,
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "type": "practice_coaching",
            })
            self._save_session(session)

            # Calculate progress percentage (step 1=0%, step 5=80%, completing step 5=100%)
            practice_progress = min(100, (practice_step - 1) * 20)
            if practice_step == 5 and session.get("qa_questions_asked", 0) >= 3:
                practice_progress = 100

            return {
                "success": True,
                "current_topic": session["current_topic"],
                "turn_number": student_turns,
                "hints_used": session.get("hints_given", 0),
                "session_id": session_id,
                "ai_response": ai_response,
                "practice_step": practice_step,
                "practice_step_name": self.PRACTICE_STEPS.get(practice_step, {}).get("name", ""),
                "practice_progress": practice_progress,
                "suggested_questions": suggested_questions,
            }

        else:
            # Main: record only — no AI response
            self._save_session(session)
            return {
                "success": True,
                "current_topic": session["current_topic"],
                "turn_number": student_turns,
                "hints_used": session.get("hints_given", 0),
                "session_id": session_id,
            }

    def seminar_interaction(
        self,
        session_id: str,
        interaction_type: str,
        student_message: str = "",
    ) -> Dict[str, Any]:
        """
        Handle special interactions: repeat, clarify, not_understood.
        """
        session = self._load_session(session_id)
        if not session:
            return {"error": "Session not found", "success": False}
        if session.get("status") == "ended":
            return {"error": "Session already ended", "success": False}

        # Find last AI message
        last_ai_msg = ""
        for msg in reversed(session["messages"]):
            if msg.get("role") == "ai":
                last_ai_msg = msg.get("content", "")
                break

        if interaction_type == "repeat":
            prompt = f"""The student asked you to repeat your last point. Rephrase this in simpler terms:

Your last message: "{last_ai_msg}"

Rephrase it more simply, using everyday examples. Start with "Sure! Let me explain again..."
Keep it under 3 sentences."""

        elif interaction_type == "clarify":
            prompt = f"""The student needs clarification. They said: "{student_message}"

Your last message was: "{last_ai_msg}"

Provide a clearer explanation, use an analogy or example. Start with "Great question! Think of it this way..."
Keep it under 4 sentences."""

        elif interaction_type == "not_understood":
            prompt = f"""The student said they didn't understand. Your last message was:
"{last_ai_msg}"

Explain the same concept in the SIMPLEST possible way, as if explaining to someone who knows nothing about this topic.
Use a real-world analogy. Start with "No worries! Let me break it down..."
Keep it under 4 sentences."""

        else:
            return {"error": f"Unknown interaction type: {interaction_type}", "success": False}

        system_prompt = self._seminar_system_prompt(
            session["subject"], session["unit_number"],
            session["current_topic"], session["candidate_name"],
            session.get("rag_context", ""),
        )

        ai_response = self._call_llm([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ])

        # Record interaction
        if student_message:
            session["messages"].append({
                "role": "student",
                "content": student_message,
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "type": interaction_type,
            })

        session["messages"].append({
            "role": "ai",
            "content": ai_response,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "type": f"response_to_{interaction_type}",
        })
        self._save_session(session)

        return {
            "success": True,
            "ai_response": ai_response,
            "interaction_type": interaction_type,
            "session_id": session_id,
        }

    def end_seminar(self, session_id: str, uploaded_content: Optional[str] = None) -> Dict[str, Any]:
        """
        End the seminar session.

        Behavior varies by session_mode:
          - "main":     Score the student + generate a response_message + track performance.
          - "demo":     Score the student but skip performance tracking.
          - "practice": No scoring — just confirm deletion of temp session file.
        """
        session = self._load_session(session_id)
        if not session:
            return {"error": "Session not found", "success": False}
        if session.get("status") == "ended":
            return {"error": "Session already ended", "success": False}

        session_mode = session.get("session_mode", "main")

        # ── Practice mode: no scoring, delete temp file ───────────────────
        if session_mode == "practice":
            session["status"] = "ended"
            session["ended_at"] = datetime.now(timezone.utc).isoformat()

            # Delete the temp session file from practice_session/
            session_path = self._session_path(session_id, "practice")
            try:
                if session_path.exists():
                    session_path.unlink()
                    print(f"  🗑️ [SeminarEngine] Deleted practice session file: {session_path.name}")
            except OSError as e:
                print(f"  ⚠️ [SeminarEngine] Failed to delete practice session file: {e}")

            # Auto-delete seminar chat history
            self.end_seminar_chat(session_id)

            return {
                "success": True,
                "session_id": session_id,
                "session_mode": "practice",
                "message": "Practice session ended. Session data has been deleted.",
            }

        # ── Score the seminar (main + demo) ───────────────────────────────
        conversation = ""
        for m in session.get("messages", []):
            role_label = "Student" if m.get("role") == "student" else "Examiner"
            conversation += f"\n{role_label}: {m.get('content', '')}\n"

        file_context = f"\n- Student Uploaded Presentation File:\n{uploaded_content[:5000]}" if uploaded_content else ""

        scoring_prompt = f"""Evaluate this student seminar session on "{session['current_topic']}" from {session['subject']}.

## CONVERSATION
{conversation[:10000]}

## CONTEXT
- Topic covered: {session['current_topic']}
- Hints given: {session.get('hints_given', 0)}
- Total student interactions: {session.get('total_interactions', 0)}{file_context}

## SCORING CRITERIA
- conceptual_understanding (0-30): Accuracy of explanations vs textbook content
- depth_of_knowledge (0-25): Detail level, examples provided, connections made
- presentation_flow (0-20): Logical sequence, clear transitions, structured explanation
- engagement (0-15): Response to AI prompts, willingness to elaborate, questions asked
- hints_penalty (0-10): 10 = no hints used, 0 = heavily relied on hints. Deduct 1 point per hint.

Return JSON:
{{"conceptual_understanding": 25, "depth_of_knowledge": 20, "presentation_flow": 16, "engagement": 12, "hints_penalty": {max(0, 10 - session.get('hints_given', 0))}, "total_score": 83, "overall_feedback": "...", "strengths": ["..."], "improvements": ["..."], "topics_mastered": ["..."], "topics_need_work": ["..."]}}"""

        scores = self._call_llm_json(
            [{"role": "user", "content": scoring_prompt}], temperature=0.2
        )

        if not scores:
            scores = {
                "conceptual_understanding": 20,
                "depth_of_knowledge": 15,
                "presentation_flow": 12,
                "engagement": 10,
                "hints_penalty": max(0, 10 - session.get("hints_given", 0)),
                "total_score": 57 + max(0, 10 - session.get("hints_given", 0)),
                "overall_feedback": "Seminar session completed.",
                "strengths": ["Participated in the seminar"],
                "improvements": ["Practice explaining concepts in more detail"],
                "topics_mastered": [],
                "topics_need_work": [session["current_topic"]],
            }

        # Ensure total is correct
        scores["total_score"] = (
            scores.get("conceptual_understanding", 0)
            + scores.get("depth_of_knowledge", 0)
            + scores.get("presentation_flow", 0)
            + scores.get("engagement", 0)
            + scores.get("hints_penalty", 0)
        )

        session["status"] = "ended"
        session["scores"] = scores
        session["ended_at"] = datetime.now(timezone.utc).isoformat()

        self._save_session(session)

        # ── Generate response_message for MAIN sessions ───────────────────
        response_message = None
        if session_mode == "main":
            feedback_prompt = f"""You just finished examining a student's seminar on "{session['current_topic']}" from {session['subject']}.

Their score: {scores['total_score']}/100
Strengths: {scores.get('strengths', [])}
Areas to improve: {scores.get('improvements', [])}

Give a brief, warm, encouraging closing message (3-4 sentences):
1. Thank them for their presentation.
2. Highlight one strength.
3. Suggest one specific area to improve.
4. End with encouragement for next time.

Keep it conversational and supportive."""

            response_message = self._call_llm([
                {"role": "system", "content": "You are a warm, encouraging seminar examiner giving closing feedback to a student."},
                {"role": "user", "content": feedback_prompt},
            ])

        # ── Update student performance (MAIN sessions only) ────────────────
        if session_mode == "main":
            try:
                from student_performance import get_performance_tracker
                tracker = get_performance_tracker()

                section_scores = {session["current_topic"]: float(scores["total_score"])}

                tracker.record_seminar_score(
                    candidate_id=session["candidate_id"],
                    subject=session["subject"],
                    unit_number=session["unit_number"],
                    section_scores=section_scores,
                    total_score=float(scores["total_score"]),
                    points=max(0, scores["total_score"]),
                    unit_title=session.get("unit_name", ""),
                    candidate_name=session.get("candidate_name", ""),
                )

                tracker.record_tutor_interaction(
                    candidate_id=session["candidate_id"],
                    subject=session["subject"],
                    unit_number=session["unit_number"],
                    topic=session["current_topic"],
                    query_summary=f"Seminar completed. Score: {scores['total_score']}/100",
                    candidate_name=session.get("candidate_name", ""),
                    unit_title=session.get("unit_name", ""),
                )
            except Exception as e:
                print(f"  ⚠️ [SeminarEngine] Failed to update performance: {e}")
        else:
            print(f"  ℹ️ [SeminarEngine] Demo session — skipping performance tracking")

        # ── Auto-delete seminar chat history ─────────────────────────────────
        self.end_seminar_chat(session_id)

        result = {
            "success": True,
            "session_id": session_id,
            "scores": scores,
            "topics_covered": [session["current_topic"]],
            "hints_used": session.get("hints_given", 0),
        }

        # Include response_message for main sessions
        if response_message:
            result["response_message"] = response_message

        return result

    # ── Demo/Practice Guidance ─────────────────────────────────────────────

    def guide_student(self, session_id: str) -> Dict[str, Any]:
        """
        Provide RAG-based guidance to the student during a DEMO session.

        Reads the full conversation history from the session,
        analyzes what the student has covered vs what the textbook expects,
        and returns structured coaching advice.

        Only available for session_mode == "demo".
        """
        session = self._load_session(session_id)
        if not session:
            return {"error": "Session not found", "success": False}
        if session.get("status") == "ended":
            return {"error": "Session already ended", "success": False}
        if session.get("session_mode") != "demo":
            return {
                "error": "Guidance is only available for demo/practice sessions",
                "success": False,
            }

        # ── Build conversation summary from session history ────────────────
        conversation = ""
        student_points = []
        for m in session.get("messages", []):
            role_label = "Student" if m.get("role") == "student" else "Coach"
            conversation += f"\n{role_label}: {m.get('content', '')}\n"
            if m.get("role") == "student":
                student_points.append(m.get("content", ""))

        if not student_points:
            # Student hasn't said anything yet
            return {
                "success": True,
                "guidance": (
                    "You haven't started presenting yet! "
                    "Take a deep breath and begin by defining the topic. "
                    "What is the first thing you'd say to introduce this subject?"
                ),
                "session_id": session_id,
                "type": "guidance",
            }

        # ── Generate structured guidance using RAG context ─────────────────
        rag_context = session.get("rag_context", "")
        guidance_prompt = f"""You are a seminar coach helping a student practice their presentation on "{session['current_topic']}" from {session['subject']}, Unit {session['unit_number']}.

## TEXTBOOK CONTENT (what the student SHOULD cover)
{rag_context if rag_context else "(No textbook context available)"}

## CONVERSATION SO FAR
{conversation[:8000]}

## YOUR TASK
Analyze the student's presentation so far and provide structured guidance:

1. **What they covered well**: List specific points the student explained correctly (be specific, reference their actual words).
2. **What's missing**: Based on the textbook content, what important points has the student NOT yet covered?
3. **What to say next**: Give them a concrete suggestion for what to present next — phrase it as "Now try explaining..." or "Next, you should talk about..."
4. **Presentation tips**: One specific tip to improve how they're presenting (e.g., "Use an example", "Define the term first", "Connect it to the previous point").

Keep the total response under 8 sentences. Be encouraging but specific."""

        system_prompt = self._demo_system_prompt(
            session["subject"], session["unit_number"],
            session["current_topic"], session["candidate_name"],
            rag_context,
        )

        guidance = self._call_llm([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": guidance_prompt},
        ])

        # Record guidance in session messages
        session["messages"].append({
            "role": "ai",
            "content": guidance,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "type": "guidance",
        })
        self._save_session(session)

        student_turns = sum(1 for m in session["messages"] if m.get("role") == "student")

        return {
            "success": True,
            "guidance": guidance,
            "session_id": session_id,
            "current_topic": session["current_topic"],
            "student_turns_so_far": student_turns,
            "type": "guidance",
        }

    # ── Helper Methods ────────────────────────────────────────────────────────

    def _give_hint(self, session: Dict[str, Any]) -> Dict[str, Any]:
        """Generate a contextual hint based on the last conversation point."""
        last_exchange = ""
        for msg in reversed(session["messages"]):
            last_exchange = f"{msg.get('role', '')}: {msg.get('content', '')}\n" + last_exchange
            if len(last_exchange) > 600:
                break

        hint_prompt = f"""The student has been silent for a while during a seminar on "{session['current_topic']}".

Last conversation:
{last_exchange}

Give a SHORT, helpful hint (1-2 sentences) to help them continue their explanation.
Don't give away the answer — just nudge them in the right direction.
Start with "💡 Here's a hint:" """

        system_prompt = self._seminar_system_prompt(
            session["subject"], session["unit_number"],
            session["current_topic"], session["candidate_name"],
            session.get("rag_context", ""),
        )

        hint = self._call_llm([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": hint_prompt},
        ])

        session["hints_given"] = session.get("hints_given", 0) + 1
        session["messages"].append({
            "role": "ai",
            "content": hint,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "type": "hint",
        })
        self._save_session(session)

        return {
            "success": True,
            "ai_response": hint,
            "is_hint": True,
            "hints_used": session["hints_given"],
            "session_id": session["session_id"],
        }


    def get_session(self, session_id: str) -> Optional[Dict[str, Any]]:
        """Get session data (without rag_context)."""
        session = self._load_session(session_id)
        if session:
            return {k: v for k, v in session.items() if k != "rag_context"}
        return None

    def get_seminar_history(
        self,
        candidate_id: str,
        subject: Optional[str] = None,
    ) -> List[Dict[str, Any]]:
        """List all seminar sessions for a student (main, demo, and practice)."""
        history = []
        # Search main, practice_data, and practice_session directories
        for base_dir in (self.data_dir, self.practice_data_dir, self.practice_session_dir):
            for path in sorted(base_dir.glob("seminar__*.json")):
                try:
                    data = json.loads(path.read_text(encoding="utf-8"))
                    if data.get("candidate_id") != candidate_id:
                        continue
                    if subject and data.get("subject", "").lower() != subject.lower():
                        continue
                    history.append({
                        "session_id": data.get("session_id"),
                        "current_topic": data.get("current_topic"),
                        "subject": data.get("subject"),
                        "unit_number": data.get("unit_number"),
                        "session_mode": data.get("session_mode", "main"),
                        "status": data.get("status"),
                        "hints_used": data.get("hints_given", 0),
                        "total_score": data.get("scores", {}).get("total_score") if data.get("scores") else None,
                        "topics_covered": [data.get("current_topic", "")] if data.get("current_topic") else [],
                        "created_at": data.get("created_at"),
                        "ended_at": data.get("ended_at"),
                    })
                except (json.JSONDecodeError, OSError):
                    continue

        history.sort(key=lambda x: x.get("created_at", ""), reverse=True)
        return history

    def get_attended_topics(
        self,
        candidate_id: str,
        subject: Optional[str] = None,
        unit_number: Optional[int] = None,
    ) -> List[Dict[str, Any]]:
        """Get all topics this student has attended in seminar sessions (main only)."""
        attended = []
        # Only count main sessions for attended topics (not practice)
        for path in sorted(self.data_dir.glob("seminar__*.json")):
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
                if data.get("candidate_id") != candidate_id:
                    continue
                if subject and data.get("subject", "").lower() != subject.lower():
                    continue
                if unit_number is not None and data.get("unit_number") != unit_number:
                    continue
                if data.get("status") != "ended":
                    continue

                topics_covered = [data.get("current_topic", "")] if data.get("current_topic") else []
                score = data.get("scores", {}).get("total_score") if data.get("scores") else None

                for topic in topics_covered:
                    if topic:
                        attended.append({
                            "topic": topic,
                            "score": score,
                            "session_id": data.get("session_id"),
                            "status": data.get("status"),
                            "unit_number": data.get("unit_number"),
                            "subject": data.get("subject"),
                            "date": data.get("ended_at") or data.get("created_at"),
                        })
            except (json.JSONDecodeError, OSError):
                continue

        return attended

    # ── Post-Session Seminar Chat ──────────────────────────────────────────

    def _chat_path(self, session_id: str) -> Path:
        """Return path for a seminar chat file."""
        return self.chat_dir / f"chat__{session_id}.json"

    def _load_chat(self, session_id: str) -> Optional[Dict[str, Any]]:
        path = self._chat_path(session_id)
        if path.exists():
            try:
                return json.loads(path.read_text(encoding="utf-8"))
            except (json.JSONDecodeError, OSError):
                pass
        return None

    def _save_chat(self, chat: Dict[str, Any]) -> None:
        path = self._chat_path(chat["session_id"])
        chat["updated_at"] = datetime.now(timezone.utc).isoformat()
        path.write_text(
            json.dumps(chat, indent=2, ensure_ascii=False), encoding="utf-8"
        )

    def _seminar_chat_system_prompt(
        self,
        session: Dict[str, Any],
    ) -> str:
        """System prompt for the post-session seminar chat."""
        candidate_name = session.get("candidate_name", "Student")
        topic = session.get("current_topic", "")
        subject = session.get("subject", "")
        unit_number = session.get("unit_number", "")
        scores = session.get("scores", {})

        # Build conversation summary
        conversation_summary = ""
        for m in session.get("messages", [])[-20:]:
            role_label = "Student" if m.get("role") == "student" else "Examiner"
            conversation_summary += f"\n{role_label}: {m.get('content', '')}\n"

        # Include uploaded file content if available
        uploaded_section = ""
        if session.get("content_source") == "uploaded_file" and session.get("session_context"):
            uploaded_section = f"""

## STUDENT'S UPLOADED PRESENTATION MATERIAL
The student uploaded a PDF/PPT file for this session. Use this as your PRIMARY reference when answering their questions.
{session['session_context'][:6000]}
"""

        return f"""You are GradeUp AI Seminar Teacher — a warm, knowledgeable teacher assisting a student during their live seminar session.

## YOUR ROLE
You are chatting with {candidate_name} who is currently presenting a seminar on "{topic}" from {subject}, Unit {unit_number}.
This is a LIVE session chat. The student can ask you questions about the topic while they are presenting.

## RECENT CONVERSATION FROM THE SEMINAR
{conversation_summary[:5000]}
{uploaded_section}
## WHAT YOU CAN DO
1. **Answer topic questions**: Answer any questions about "{topic}" — explain concepts, provide examples, clarify doubts.
2. **Help with presentation**: If the student is unsure what to cover next, suggest key points from their uploaded material or the textbook.
3. **Clarify concepts**: If they're confused about something, explain it simply. Reference their own uploaded material when possible.
4. **Encourage**: Be positive and motivating.

## RULES
- Be warm, friendly, and educational — like a helpful teacher during class.
- Keep responses conversational and SHORT (2-4 sentences).
- Stay focused on the topic "{topic}" and the seminar.
- When the student has uploaded material, prioritize referencing THEIR content over general knowledge.
- NEVER produce inappropriate, 18+, or harmful content.
- If asked about unrelated topics, gently redirect: "Let's focus on your seminar! What would you like to know about {topic}?"
"""

    def start_seminar_chat(self, session_id: str) -> Dict[str, Any]:
        """
        Start a live session chat for an active seminar.

        For PRACTICE sessions: chat is merged into the session itself —
        returns existing session messages. No separate chat file created.

        For MAIN/DEMO sessions: creates a temporary chat file in
        seminar_chat/ folder. Auto-deleted when session ends.
        """
        session = self._load_session(session_id)
        if not session:
            return {"error": "Session not found", "success": False}
        if session.get("status") == "ended":
            return {"error": "Chat is not available after the session has ended", "success": False}

        # ── Practice mode: chat IS the session — no separate chat file ────
        if session.get("session_mode") == "practice":
            return {
                "success": True,
                "session_id": session_id,
                "message": "Practice session uses a unified endpoint. Use /seminar/respond for all interactions.",
                "ai_message": session["messages"][0]["content"] if session.get("messages") else "",
                "current_topic": session.get("current_topic"),
                "chat_messages": session.get("messages", []),
            }

        # ── Main / Demo: separate chat file ───────────────────────────────
        # Check if chat already exists
        existing = self._load_chat(session_id)
        if existing and existing.get("status") == "active":
            return {
                "success": True,
                "session_id": session_id,
                "message": "Chat session already active",
                "chat_messages": existing.get("messages", []),
            }

        # Generate welcome message
        system_prompt = self._seminar_chat_system_prompt(session)
        welcome_prompt = f"""The student {session.get('candidate_name', 'Student')} has started a live seminar session on "{session.get('current_topic', '')}".

Greet them warmly, let them know they can ask you questions about the topic anytime during their presentation, and wish them good luck.

Keep it to 2-3 sentences. Be encouraging!"""

        ai_welcome = self._call_llm([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": welcome_prompt},
        ])

        chat = {
            "session_id": session_id,
            "candidate_id": session.get("candidate_id"),
            "candidate_name": session.get("candidate_name"),
            "subject": session.get("subject"),
            "current_topic": session.get("current_topic"),
            "status": "active",
            "messages": [
                {
                    "role": "ai",
                    "content": ai_welcome,
                    "timestamp": datetime.now(timezone.utc).isoformat(),
                }
            ],
            "created_at": datetime.now(timezone.utc).isoformat(),
        }

        self._save_chat(chat)

        return {
            "success": True,
            "session_id": session_id,
            "ai_message": ai_welcome,
            "current_topic": session.get("current_topic"),
        }

    def seminar_chat_respond(
        self,
        session_id: str,
        student_message: str,
    ) -> Dict[str, Any]:
        """
        Handle a student message in the seminar chat.

        For PRACTICE sessions: delegates to seminar_respond() — single
        unified endpoint handles both session and chat.

        For MAIN/DEMO sessions: uses the separate chat file.
        """
        # ── Practice mode: delegate to seminar_respond ────────────────────
        session = self._load_session(session_id)
        if session and session.get("session_mode") == "practice":
            return self.seminar_respond(
                session_id=session_id,
                student_message=student_message,
            )

        # ── Main / Demo: use separate chat file ──────────────────────────
        chat = self._load_chat(session_id)
        if not chat:
            return {"error": "Chat session not found. Start the chat first.", "success": False}
        if chat.get("status") != "active":
            return {"error": "Chat session has ended", "success": False}

        if not session:
            return {"error": "Original seminar session not found", "success": False}

        # Record student message
        chat["messages"].append({
            "role": "student",
            "content": student_message,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })

        # Build LLM messages
        system_prompt = self._seminar_chat_system_prompt(session)
        llm_messages = [{"role": "system", "content": system_prompt}]
        for msg in chat["messages"][-16:]:
            role = "assistant" if msg["role"] == "ai" else "user"
            llm_messages.append({"role": role, "content": msg["content"]})

        ai_response = self._call_llm(llm_messages)

        chat["messages"].append({
            "role": "ai",
            "content": ai_response,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })
        self._save_chat(chat)

        return {
            "success": True,
            "ai_response": ai_response,
            "session_id": session_id,
            "message_count": len(chat["messages"]),
        }

    def end_seminar_chat(self, session_id: str) -> Dict[str, Any]:
        """
        End the post-session seminar chat and delete the chat history.
        """
        chat_path = self._chat_path(session_id)
        if not chat_path.exists():
            return {"error": "No chat session found for this session", "success": False}

        try:
            chat_path.unlink()
            print(f"  🗑️ [SeminarEngine] Deleted chat history for session {session_id}")
        except OSError as e:
            print(f"  ⚠️ [SeminarEngine] Failed to delete chat file: {e}")

        return {
            "success": True,
            "session_id": session_id,
            "message": "Chat session ended and history deleted",
        }


# ── Global singleton ──────────────────────────────────────────────────────

_seminar_engine: Optional[SeminarEngine] = None


def get_seminar_engine() -> SeminarEngine:
    global _seminar_engine
    if _seminar_engine is None:
        _seminar_engine = SeminarEngine()
    return _seminar_engine
